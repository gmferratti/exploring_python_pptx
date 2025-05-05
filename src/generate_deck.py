# src/generate_deck.py
import argparse
from pptx import Presentation
from copy import deepcopy
from pptx.util import Pt
from template_utils import create_base_template

# Mapeamento layout -> índice em layouts
LAYOUT_INDEX = {
    "cover": 0,
    "content": 1,
    "closing": 2
}

def duplicate_slide(src_slide, dest_prs):
    """
    Duplica shapes de src_slide para um novo slide em dest_prs.
    """
    blank = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank)
    for shape in src_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

def fill_placeholders(slide, data):
    """
    Preenche placeholders nomeados em cada slide duplicado.
    Chaves esperadas:
      cover: title, subtitle
      content: title, bullets (lista)
      closing: text
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name.lower()
        tf = shape.text_frame
        # Cover
        if "cover" in name and "title" in data:
            tf.text = data["title"]
            tf.paragraphs[0].font.size = Pt(36)
        if "cover" in name and "subtitle" in data:
            tf.text = data["subtitle"]
            tf.paragraphs[0].font.size = Pt(20)
        # Content
        if "content" in name and "title" in data:
            tf.text = data["title"]
            tf.paragraphs[0].font.size = Pt(28)
        if "content" in name and "bullets" in data:
            bullets = data["bullets"]
            if bullets:
                tf.text = bullets[0]
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
                    p.font.size = Pt(16)
        # Closing
        if "closing" in name and "text" in data:
            tf.text = data["text"]
            tf.paragraphs[0].font.size = Pt(36)

def build_deck(output_path, slides_data):
    """
    Gera deck final sem arquivo intermediário.
    slides_data: lista de dicts com chaves layout/title/bullets/text...
    """
    base_prs, layouts = create_base_template()
    deck = Presentation()
    # Remove slide em branco inicial
    while deck.slides:
        rId = deck.slides._sldIdLst[0].rId
        deck.part.drop_rel(rId)
        deck.slides._sldIdLst.remove(deck.slides._sldIdLst[0])
    # Duplica e preenche
    for data in slides_data:
        idx = LAYOUT_INDEX.get(data["layout"], 1)
        new_slide = duplicate_slide(layouts[idx], deck)
        fill_placeholders(new_slide, data)
    deck.save(output_path)
    print(f"Deck salvo em: {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "-o", "--out", default="presentations/deck_final.pptx",
        help="Caminho de saída do PPTX"
    )
    args = parser.parse_args()
    # Exemplo de slides_data; customize conforme necessário
    slides_data = [
        {
            "layout": "cover",
            "title": "ML Engineering com Kedro e MLflow",
            "subtitle": "Time de MLOps • Data Science Ipiranga"
        },
        {
            "layout": "content",
            "title": "Agenda",
            "bullets": [
                "MLOps & ML Engineering",
                "Por que Kedro?",
                "Por que MLflow?",
                "Integração Kedro + MLflow",
                "Hands‑on passo a passo",
                "Boas práticas & próximos passos"
            ]
        },
        {
            "layout": "closing",
            "text": "Obrigado! • Perguntas?"
        }
    ]
    build_deck(args.out, slides_data)