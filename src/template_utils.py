# src/template_utils.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Paleta de cores Ipiranga
BLUE = RGBColor(0x00, 0x57, 0xB6)
YELLOW = RGBColor(0xFF, 0xD5, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_branding(slide, prs, footer_text="MLOps • Data Science Ipiranga"):
    """
    Adiciona barra amarela (topo), barra azul (base) e rodapé com texto.
    """
    # Top bar
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = YELLOW
    top.line.fill.background()
    # Bottom bar
    bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35),
        prs.slide_width, Inches(0.35)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = BLUE
    bottom.line.fill.background()
    # Footer text
    footer = slide.shapes.add_textbox(
        Inches(0.2), prs.slide_height - Inches(0.3),
        Inches(5), Inches(0.3)
    )
    p = footer.text_frame.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE

def create_base_template():
    """
    Cria um Presentation com 3 slides-base (layouts):
        0: cover (com placeholders nomeados)
        1: content
        2: closing
    Retorna (prs, layouts).
    """
    prs = Presentation()
    blank = prs.slide_layouts[6]
    layouts = []

    # --- Cover ---
    slide = prs.slides.add_slide(blank)
    add_branding(slide, prs)
    # Placeholder de título
    shp_title = slide.shapes.add_textbox(Inches(0.8), Inches(2),
                                            Inches(9), Inches(1.2))
    shp_title.name = "cover_title"
    p = shp_title.text_frame.paragraphs[0]
    p.text = ""
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE
    # Placeholder de subtítulo
    shp_sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.3),
                                        Inches(9), Inches(1))
    shp_sub.name = "cover_subtitle"
    p2 = shp_sub.text_frame.paragraphs[0]
    p2.text = ""
    p2.font.size = Pt(20)
    p2.font.color.rgb = BLUE
    layouts.append(slide)

    # --- Content ---
    slide = prs.slides.add_slide(blank)
    add_branding(slide, prs)
    shp_ctitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.6),
                                            Inches(8.5), Inches(1))
    shp_ctitle.name = "content_title"
    pt = shp_ctitle.text_frame.paragraphs[0]
    pt.text = ""
    pt.font.size = Pt(28)
    pt.font.bold = True
    pt.font.color.rgb = BLUE
    shp_cbody = slide.shapes.add_textbox(Inches(0.8), Inches(1.7),
                                            Inches(8.5), Inches(4))
    shp_cbody.name = "content_body"
    # deixa o corpo vazio
    layouts.append(slide)

    # --- Closing ---
    slide = prs.slides.add_slide(blank)
    add_branding(slide, prs)
    shp_close = slide.shapes.add_textbox(Inches(1.5), Inches(2.5),
                                            Inches(7), Inches(2))
    shp_close.name = "closing_text"
    p3 = shp_close.text_frame.paragraphs[0]
    p3.text = ""
    p3.font.size = Pt(36)
    p3.font.bold = True
    p3.font.color.rgb = BLUE
    layouts.append(slide)

    return prs, layouts